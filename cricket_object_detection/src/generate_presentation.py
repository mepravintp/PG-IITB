from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Helper to add a slide with title and content
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1: Title
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Cricket Object Detection Project"
slide.placeholders[1].text = "Detecting and classifying cricket objects at the grid-cell level\nPrepared by: [Your Name]\nDate: December 7, 2025"

# Slide 2: Table of Contents
add_slide("Table of Contents", 
"""1. Project Overview
2. Folder Structure
3. Step-by-Step Workflow
4. Key Scripts and Pipelines
5. Model Training & Evaluation
6. Visualizations
7. Learnings & Challenges
8. References""")

# Slide 3: Project Overview
add_slide("Project Overview", 
"""Goal: Detect and classify cricket objects (bat, ball, stumps, no_object) at the grid-cell level in images.
Approach: Manual annotation, feature extraction, classical ML models (Random Forest, KNN), and visualization.""")

# Slide 4: Folder Structure
add_slide("Folder Structure", 
"""cricket_object_detection/
├── data/
├── models/
├── outputs/
├── src/
├── requirements.txt
├── README.md
└── TASKS.csv
Data, models, outputs, source code, and documentation are organized for clarity.""")

# Slide 5: Step-by-Step Workflow
add_slide("Step-by-Step Workflow", 
"""1. Data Collection
2. Preprocessing (resize/crop to 800x600)
3. Annotation (8x8 grid, 0: no_object, 1: ball, 2: bat, 3: stump)
4. Manual Annotation (GUI/script)
5. Feature Extraction (HOG, grayscale, etc.)
6. Save Features (CSV)
7. Model Training (Random Forest, KNN)
8. Evaluation (metrics, visualizations)
9. Prediction & Visualization (overlay grid, save results)""")

# Slide 6: Annotation Example
add_slide("Annotation Example", 
"Screenshot or diagram of annotation tool/grid overlay\nExplain cell tagging process")

# Slide 7: Key Scripts and Pipelines
add_slide("Key Scripts and Pipelines", 
"""preprocess.py: Image resizing/cropping
annotate.py: Annotation utilities
annotation_to_csv.py: Convert annotations to CSV
extract_features.py: Feature extraction
cell_level_pipeline.py: Full pipeline""")

# Slide 8: Model Training & Evaluation
add_slide("Model Training & Evaluation", 
"""Trained Random Forest and KNN classifiers
Compared using accuracy, precision, recall, F1
Show classification report/bar graph (insert image)""")

# Slide 9: Visualizations
add_slide("Visualizations", 
"""Overlaid predicted cell labels on images
Color codes: gray (no_object), red (ball), green (bat), blue (stump)
Show example images (insert images)""")

# Slide 10: Learnings & Challenges
add_slide("Learnings & Challenges", 
"""Importance of consistent annotation and preprocessing
Feature engineering is effective for classical ML
Cell-level features are crucial for fine-grained classification
Model comparison highlights trade-offs
Visualization is key for debugging""")

# Slide 11: References
add_slide("References", 
"""Project scripts and notebooks
latest_readme.md
scikit-learn, scikit-image documentation""")

# Save the presentation
prs.save("cricket_object_detection_presentation.pptx")
print("Presentation saved as cricket_object_detection_presentation.pptx")
